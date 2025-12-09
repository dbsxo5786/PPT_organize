from pptx import Presentation
import logging

log = logging.getLogger(__name__)

def parse_ppt(file_path: str) -> dict[str, str]:
    prs = Presentation(file_path)
    slides_text = {}

    for i, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            # Shape에 text 속성 있으면 리스트에 추가
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())
                
        slides_text[str(i)] = " ".join(texts)

    log.info("Text in slides: %s", slides_text)
    return slides_text